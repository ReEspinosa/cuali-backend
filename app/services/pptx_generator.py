"""
Genera un archivo .pptx con diseño infantil estilo Cuali:
fondo de color sólido, círculos blancos semitransparentes (efecto liquid
glass), emoji grande como ilustración, y tipografía grande y legible.

Los temas de color disponibles son: azul, rosa, amarillo, verde, naranja.
"""

import os
import random

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUTPUT_DIR = "generated_docs"

# Fondo sólido (tono medio/oscuro para que el texto blanco sea legible)
TEMAS = {
    "azul": RGBColor(0x4E, 0x7A, 0xB5),
    "rosa": RGBColor(0xD4, 0x53, 0x7E),
    "amarillo": RGBColor(0xBA, 0x75, 0x17),
    "verde": RGBColor(0x63, 0x99, 0x22),
    "naranja": RGBColor(0xD8, 0x5A, 0x30),
}

FUENTE = "Comic Sans MS"  # legible y amigable para primaria; disponible en Windows/Mac/Office

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_alpha(shape, alpha_pct: int) -> None:
    """Aplica transparencia a un relleno sólido (alpha_pct: 0-100, donde 100 = opaco)."""
    solid_fill = shape.fill._xPr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(alpha_pct * 1000)})
    srgb.append(alpha)


def _fondo(slide, color: RGBColor) -> None:
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = color
    fondo.line.fill.background()
    fondo.shadow.inherit = False


def _circulos_blancos(slide, cantidad: int = 6, semilla: int = 0) -> None:
    rng = random.Random(semilla)  # semilla fija por diapositiva: mismo diseño en cada regeneración
    for _ in range(cantidad):
        tam = Emu(int(Inches(rng.uniform(1.5, 4.0))))
        x = Emu(int(Inches(rng.uniform(-1.5, 12.5))))
        y = Emu(int(Inches(rng.uniform(-1.5, 6.5))))
        circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, tam, tam)
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_alpha(circulo, rng.randint(10, 22))
        circulo.line.fill.background()
        circulo.shadow.inherit = False


def _texto(slide, texto: str, x, y, w, h, size: int, bold: bool = False,
           align=PP_ALIGN.LEFT) -> None:
    caja = slide.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.name = FUENTE
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _bullets(slide, puntos: list[str], x, y, w, h, size: int) -> None:
    caja = slide.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    for i, punto in enumerate(puntos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = f"•  {punto}"
        run.font.name = FUENTE
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def generar_pptx_diapositivas(recurso_id: str, datos: dict) -> str:
    """
    datos = {
      "titulo": str,
      "tema_color": "azul" | "rosa" | "amarillo" | "verde" | "naranja",
      "diapositivas": [ { "titulo": str, "puntos": [str], "emoji": str } ]
    }
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, f"diapositivas_{recurso_id}.pptx")

    color = TEMAS.get(datos.get("tema_color", "azul"), TEMAS["azul"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # --- Portada ---
    slide = prs.slides.add_slide(blank)
    _fondo(slide, color)
    _circulos_blancos(slide, cantidad=8, semilla=1)
    primera = datos.get("diapositivas", [{}])[0] if datos.get("diapositivas") else {}
    emoji_portada = primera.get("emoji", "⭐")
    _texto(slide, emoji_portada, Inches(0), Inches(1.2), SLIDE_W, Inches(1.6), 88, align=PP_ALIGN.CENTER)
    _texto(slide, datos.get("titulo", "Mi presentación"), Inches(1), Inches(3.1), Inches(11.3), Inches(2.2),
           48, bold=True, align=PP_ALIGN.CENTER)

    # --- Diapositivas de contenido ---
    for i, d in enumerate(datos.get("diapositivas", [])):
        slide = prs.slides.add_slide(blank)
        _fondo(slide, color)
        _circulos_blancos(slide, cantidad=6, semilla=i + 2)

        _texto(slide, d.get("titulo", ""), Inches(0.8), Inches(0.5), Inches(9.5), Inches(1.3), 36, bold=True)
        if d.get("emoji"):
            _texto(slide, d["emoji"], Inches(10.5), Inches(0.3), Inches(2.3), Inches(1.6), 66, align=PP_ALIGN.CENTER)

        _bullets(slide, d.get("puntos", []), Inches(1.0), Inches(2.0), Inches(11.3), Inches(5.0), 24)

    prs.save(path)
    return path